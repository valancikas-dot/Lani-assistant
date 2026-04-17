"""
pdf_tool.py – Document reading and extraction tools.

Reads text from:
  • PDF files (via pypdf / pdfminer)
  • Word documents (.docx via python-docx)
  • Excel spreadsheets (.xlsx via openpyxl)
  • PowerPoint (.pptx via python-pptx)
  • Plain text, Markdown, CSV

Does NOT require these packages to be pre-installed – gracefully degrades
to a subprocess-based fallback (pdftotext) or returns an install hint.

Tools:
  read_pdf           – extract text from PDF (page range support)
  read_docx          – extract text from Word .docx
  read_xlsx          – extract data from Excel .xlsx as CSV/JSON
  read_any_document  – auto-detect file type and extract text
  summarize_document_file – read + LLM summary in one step
"""

from __future__ import annotations

import asyncio
import csv
import io
import logging
import shutil
from pathlib import Path
from typing import Any, Dict, Optional, cast

from app.schemas.commands import ToolResult
from app.tools.base import BaseTool

log = logging.getLogger(__name__)

_MAX_CHARS = 30_000


def _cap(text: str, max_chars: int = _MAX_CHARS) -> str:
    if len(text) > max_chars:
        return text[:max_chars] + f"\n\n… [teikta {max_chars} iš {len(text)} simbolių]"
    return text


# ─────────────────────────────────────────────────────────────────────────────

async def _extract_pdf(path: Path, start_page: int = 1, end_page: int | None = None) -> str:
    """Extract text from PDF. Tries pypdf → pdfminer → pdftotext fallback."""
    # Method 1: pypdf (pure Python, fast)
    try:
        import pypdf  # type: ignore
        reader = pypdf.PdfReader(str(path))
        pages  = reader.pages
        total  = len(pages)
        start  = max(0, start_page - 1)
        end    = min(total, end_page) if end_page else total
        texts  = []
        for i in range(start, end):
            try:
                texts.append(pages[i].extract_text() or "")
            except Exception:
                texts.append("")
        return "\n\n".join(texts)
    except ImportError:
        pass

    # Method 2: pdfminer.six
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract  # type: ignore
        loop = asyncio.get_event_loop()
        text = await loop.run_in_executor(None, lambda: pdfminer_extract(str(path)))
        return text or ""
    except ImportError:
        pass

    # Method 3: pdftotext CLI (often pre-installed on macOS via poppler)
    if shutil.which("pdftotext"):
        proc = await asyncio.create_subprocess_exec(
            "pdftotext", str(path), "-",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        out, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
        return out.decode(errors="replace")

    raise RuntimeError(
        "PDF skaitymas nepasiekiamas. Įdiek: pip install pypdf\n"
        "arba: brew install poppler"
    )


async def _extract_docx(path: Path) -> str:
    """Extract text from .docx using python-docx."""
    try:
        import docx  # type: ignore
        loop = asyncio.get_event_loop()
        def _do():
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return await loop.run_in_executor(None, _do)
    except ImportError:
        raise RuntimeError("python-docx nėra. Įdiek: pip install python-docx")


async def _extract_xlsx(path: Path, as_json: bool = False) -> str:
    """Extract data from .xlsx. Returns CSV string or JSON."""
    try:
        import openpyxl  # type: ignore
        loop = asyncio.get_event_loop()
        def _do():
            wb = openpyxl.load_workbook(str(path), data_only=True)
            parts: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                buf = io.StringIO()
                writer = csv.writer(buf)
                writer.writerows(rows)
                parts.append(f"## Sheet: {sheet_name}\n{buf.getvalue()}")
            return "\n\n".join(parts)
        return await loop.run_in_executor(None, _do)
    except ImportError:
        raise RuntimeError("openpyxl nėra. Įdiek: pip install openpyxl")


async def _extract_pptx(path: Path) -> str:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
        loop = asyncio.get_event_loop()
        def _do():
            prs = Presentation(str(path))
            slides_text: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    shape_text = cast(str, getattr(shape, "text", ""))
                    if shape_text.strip():
                        texts.append(shape_text.strip())
                slides_text.append(f"### Slide {i}\n" + "\n".join(texts))
            return "\n\n".join(slides_text)
        return await loop.run_in_executor(None, _do)
    except ImportError:
        raise RuntimeError("python-pptx nėra. Įdiek: pip install python-pptx")


# ─────────────────────────────────────────────────────────────────────────────

class ReadPdfTool(BaseTool):
    name = "read_pdf"
    description = (
        "Extract text from a PDF file. Supports page range. "
        "Parameters: file_path (required), start_page (default 1), end_page (optional), "
        "max_chars (default 30000)."
    )
    requires_approval = False
    parameters = [
        {"name": "file_path",  "description": "Absolute path to PDF file", "required": True},
        {"name": "start_page", "description": "Start page (default 1)", "required": False},
        {"name": "end_page",   "description": "End page (optional)", "required": False},
        {"name": "max_chars",  "description": "Maximum characters to return", "required": False},
    ]

    async def run(self, params: Dict[str, Any]) -> ToolResult:
        file_path: str = params.get("file_path", "").strip()
        if not file_path:
            return ToolResult(tool_name=self.name, status="error", message="file_path is required")

        src = Path(file_path).expanduser()
        if not src.exists():
            return ToolResult(tool_name=self.name, status="error",
                              message=f"Failas nerastas: {file_path}")

        start = int(params.get("start_page", 1))
        end_raw = params.get("end_page")
        end = int(str(end_raw)) if end_raw not in (None, "") else None
        max_c = int(params.get("max_chars", _MAX_CHARS))

        try:
            text = await _extract_pdf(src, start, end)
        except RuntimeError as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))
        except Exception as e:
            log.exception("[read_pdf] error")
            return ToolResult(tool_name=self.name, status="error", message=str(e))

        return ToolResult(
            tool_name=self.name, status="success",
            message=f"✅ {len(text)} simboliai iš PDF: {src.name}",
            data={"text": _cap(text, max_c), "total_chars": len(text), "file": str(src)},
        )


class ReadDocxTool(BaseTool):
    name = "read_docx"
    description = (
        "Extract text from a Word .docx file. "
        "Parameters: file_path (required), max_chars (default 30000)."
    )
    requires_approval = False
    parameters = [
        {"name": "file_path", "description": "Absolute path to .docx file", "required": True},
        {"name": "max_chars", "description": "Max characters to return",    "required": False},
    ]

    async def run(self, params: Dict[str, Any]) -> ToolResult:
        file_path: str = params.get("file_path", "").strip()
        if not file_path:
            return ToolResult(tool_name=self.name, status="error", message="file_path is required")

        src = Path(file_path).expanduser()
        if not src.exists():
            return ToolResult(tool_name=self.name, status="error",
                              message=f"Failas nerastas: {file_path}")

        max_c = int(params.get("max_chars", _MAX_CHARS))
        try:
            text = await _extract_docx(src)
        except RuntimeError as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))
        except Exception as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))

        return ToolResult(
            tool_name=self.name, status="success",
            message=f"✅ {len(text)} simboliai iš: {src.name}",
            data={"text": _cap(text, max_c), "total_chars": len(text), "file": str(src)},
        )


class ReadXlsxTool(BaseTool):
    name = "read_xlsx"
    description = (
        "Read data from an Excel .xlsx spreadsheet. Returns as CSV text. "
        "Parameters: file_path (required), max_chars (default 30000)."
    )
    requires_approval = False
    parameters = [
        {"name": "file_path", "description": "Absolute path to .xlsx file", "required": True},
        {"name": "max_chars", "description": "Max characters to return",    "required": False},
    ]

    async def run(self, params: Dict[str, Any]) -> ToolResult:
        file_path: str = params.get("file_path", "").strip()
        if not file_path:
            return ToolResult(tool_name=self.name, status="error", message="file_path is required")

        src = Path(file_path).expanduser()
        if not src.exists():
            return ToolResult(tool_name=self.name, status="error",
                              message=f"Failas nerastas: {file_path}")

        max_c = int(params.get("max_chars", _MAX_CHARS))
        try:
            text = await _extract_xlsx(src)
        except RuntimeError as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))
        except Exception as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))

        return ToolResult(
            tool_name=self.name, status="success",
            message=f"✅ {len(text)} simboliai iš Excel: {src.name}",
            data={"text": _cap(text, max_c), "total_chars": len(text), "file": str(src)},
        )


class ReadAnyDocumentTool(BaseTool):
    name = "read_any_document"
    description = (
        "Auto-detect file type and extract text from PDF, DOCX, XLSX, PPTX, TXT, MD, CSV. "
        "Parameters: file_path (required), max_chars (default 30000)."
    )
    requires_approval = False
    parameters = [
        {"name": "file_path", "description": "Absolute path to document", "required": True},
        {"name": "max_chars", "description": "Max characters to return",  "required": False},
    ]

    async def run(self, params: Dict[str, Any]) -> ToolResult:
        file_path: str = params.get("file_path", "").strip()
        if not file_path:
            return ToolResult(tool_name=self.name, status="error", message="file_path is required")

        src = Path(file_path).expanduser()
        if not src.exists():
            return ToolResult(tool_name=self.name, status="error",
                              message=f"Failas nerastas: {file_path}")

        max_c = int(params.get("max_chars", _MAX_CHARS))
        ext = src.suffix.lower()

        try:
            if ext == ".pdf":
                text = await _extract_pdf(src)
            elif ext == ".docx":
                text = await _extract_docx(src)
            elif ext in (".xlsx", ".xls"):
                text = await _extract_xlsx(src)
            elif ext == ".pptx":
                text = await _extract_pptx(src)
            elif ext in (".txt", ".md", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".csv", ".html", ".xml"):
                text = src.read_text(encoding="utf-8", errors="replace")
            else:
                # Try reading as text anyway
                try:
                    text = src.read_text(encoding="utf-8", errors="replace")
                except Exception:
                    return ToolResult(tool_name=self.name, status="error",
                                      message=f"Nepalaikomas failo formatas: {ext}")
        except RuntimeError as e:
            return ToolResult(tool_name=self.name, status="error", message=str(e))
        except Exception as e:
            log.exception("[read_any_document] error")
            return ToolResult(tool_name=self.name, status="error", message=str(e))

        return ToolResult(
            tool_name=self.name, status="success",
            message=f"✅ {len(text)} simboliai iš {ext} failo: {src.name}",
            data={"text": _cap(text, max_c), "total_chars": len(text),
                  "file": str(src), "type": ext},
        )


class SummarizeDocumentFileTool(BaseTool):
    name = "summarize_document_file"
    description = (
        "Read a document file and return an AI-generated summary. "
        "Parameters: file_path (required), language ('en'|'lt', default 'en'), "
        "focus (optional: 'key_points'|'action_items'|'technical'|'general'), "
        "max_length (target summary length in words, default 300)."
    )
    requires_approval = False
    parameters = [
        {"name": "file_path",   "description": "Absolute path to document", "required": True},
        {"name": "language",    "description": "'en' or 'lt'", "required": False},
        {"name": "focus",       "description": "Summary focus", "required": False},
        {"name": "max_length",  "description": "Target word count", "required": False},
    ]

    async def run(self, params: Dict[str, Any]) -> ToolResult:
        # Step 1: read
        reader = ReadAnyDocumentTool()
        read_result = await reader.run({"file_path": params.get("file_path", ""), "max_chars": 15000})
        if read_result.status == "error":
            return read_result

        text: str   = (read_result.data or {}).get("text", "")
        language    = params.get("language", "en")
        focus       = params.get("focus", "general")
        max_length  = int(params.get("max_length", 300))

        focus_map = {
            "key_points":    "Extract the main key points as a bullet list.",
            "action_items":  "Extract all action items and tasks as a numbered list.",
            "technical":     "Summarize the technical details, specs, and implementation notes.",
            "general":       "Write a comprehensive but concise summary.",
        }

        lang_note = "Reply in Lithuanian (lietuvių kalba)." if language == "lt" else "Reply in English."
        system = f"You are a document analyst. {lang_note} {focus_map.get(focus, focus_map['general'])} Target length: ~{max_length} words."

        from app.services.llm_text_service import complete_text
        from app.core.config import settings as cfg

        if not getattr(cfg, "OPENAI_API_KEY", ""):
            return ToolResult(tool_name=self.name, status="error",
                              message="OPENAI_API_KEY nenustatytas")

        summary = await complete_text(
            openai_api_key=cfg.OPENAI_API_KEY,
            openai_model=getattr(cfg, "LLM_MODEL", "gpt-4o-mini"),
            openai_messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": f"Document content:\n\n{text}"},
            ],
            max_tokens=max_length * 2,
            temperature=0.3,
        )

        return ToolResult(
            tool_name=self.name, status="success",
            message=f"✅ Santrauka sukurta ({len(summary.split())} žodžiai)",
            data={
                "summary":   summary,
                "file":      params.get("file_path"),
                "language":  language,
                "word_count": len(summary.split()),
            },
        )
